import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Cm
from io import BytesIO
from PIL import Image
import tempfile
import os
import copy

st.title("📊 Excel + 图片生成 PowerPoint")

 # ppt_file = st.file_uploader("上传 PPT 模板 (.pptx)", type=["pptx"])
ppt_file = "1.pptx"  # 固定使用项目中的模板文件
excel_files = st.file_uploader("上传一个或多个 Excel 文件", type=["xlsx", "xls"], accept_multiple_files=True)
image_files = st.file_uploader("上传产品图片（可多选）", type=["jpg", "jpeg", "png", "bmp", "gif"], accept_multiple_files=True)

if os.path.exists(ppt_file) and excel_files:
    if st.button("生成 PowerPoint"):
        with tempfile.TemporaryDirectory() as tmpdir:
            image_folder = os.path.join(tmpdir, "images")
            os.makedirs(image_folder, exist_ok=True)
            for file in image_files:
                with open(os.path.join(image_folder, file.name), "wb") as f:
                    f.write(file.getbuffer())

            df_all = pd.concat([pd.read_excel(f) for f in excel_files], ignore_index=True)

            ppt = Presentation(ppt_file)
            template_slide = ppt.slides[0]
            new_ppt = Presentation()
            new_ppt.slide_width = ppt.slide_width
            new_ppt.slide_height = ppt.slide_height

            def clone_slide(pres, slide):
                layout = pres.slide_layouts[6]
                new_slide = pres.slides.add_slide(layout)
                for shape in slide.shapes:
                    new_shape = copy.deepcopy(shape.element)
                    new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
                return new_slide

            def replace_text(slide, data):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                text = run.text
                                for col, val in data.items():
                                    if pd.isna(val): 
                                        val = ""
                                    # 处理 ITEM# 中的 .0 问题
                                    if isinstance(val, float) and val.is_integer():
                                        val = int(val)
                                    text = text.replace(f"{{{col.strip()}}}", str(val))
                                run.text = text

            def insert_images(slide, item_number):
                matches = [f for f in os.listdir(image_folder) if f.startswith(str(item_number)) and f.lower().endswith(('.jpg','.jpeg','.png','.bmp','.gif'))]
                if not matches: return
                slide_width = new_ppt.slide_width
                slide_height = new_ppt.slide_height
                pic_height = Cm(18)
                right_margin = Cm(3)

                total_pic_height = pic_height * len(matches)
                space = (slide_height - total_pic_height) / (len(matches)+1) if len(matches) > 1 else (slide_height - total_pic_height)/2
                current_top = space

                for fname in sorted(matches):
                    img_path = os.path.join(image_folder, fname)
                    with Image.open(img_path) as img:
                        w, h = img.size
                        aspect_ratio = w / h if h else 1
                        pic_width = pic_height * aspect_ratio
                    slide.shapes.add_picture(img_path, slide_width - right_margin - pic_width, current_top, height=pic_height)
                    current_top += pic_height + space

            for _, row in df_all.iterrows():
                row_data = {k.strip(): v for k, v in row.to_dict().items()}
                for k in row_data:
                    v = row_data[k]
                    if isinstance(v, float) and v.is_integer():
                        row_data[k] = str(int(v))
                    else:
                        row_data[k] = str(v).strip()
                for field in ["Item Width (inch)", "Item Depth (inch)", "Item Height (inch)"]:
                    val = row_data.get(field, "")
                    try:
                        row_data[field] = round(float(val) * 2.54, 1) if val and not pd.isna(val) else ""
                    except:
                        row_data[field] = ""
                slide = clone_slide(new_ppt, template_slide)
                replace_text(slide, row_data)
                item_no = row_data.get("ITEM#", "")
                insert_images(slide, item_no)

            output = BytesIO()
            new_ppt.save(output)
            st.success("✅ 幻灯片已生成")
            st.download_button("📥 点击下载 PPT", output.getvalue(), file_name="output.pptx")
# 说明直接展示，不使用折叠框
st.markdown("### 📖 使用说明 / How to Use")
st.markdown("""
#### 🧾 Excel 文件要求  
- Excel文件请从SOL导出, 必须包含以下列（列名必须一致,请注意空格和大小写）:
  `ITEM#`, `Item Description`, `Item Width (inch)`, `Item Height (inch)`, `Item Depth (inch)`, `FOB NB`, `Retail AUD`  
- 尺寸单位为英寸，程序将自动转换为厘米
- 多个Excel文件可以同时处理，列顺序不限
- 可以包含其他多余的列, 程序不会处理其他多余的列.

#### 🖼 图片命名规则  
- 图片需手动命名,或者从SOL导出, 图片用5位数的Item Number命名.
  - 图片从SOL导出的要求如下
    - SOL里面确实已经插入了高清图片(用于做PPT的高质量图片)
    - SOL的第一列的内容为Item Number, 这样从SOL导出的时候多维表格会自动用Item Number来给图片命名.
- 命名格式示例： `12345.jpg`, `12345(1).jpg`, `12345(2).jpg`  
- 每个产品编号可对应多张图  
- 图像将右对齐，自动垂直分布在幻灯片上，高度固定为14cm

#### 导出Excel和图片的方法如下
""")
st.image("SOL Export.png")
st.markdown("""
    1. 制作完成以后还需要手动把PPT应用到以下的template当中去.
        - 请点击下面按钮下载真正的Kmart Buy Trip的模板
        - 下载好了以后打开, 手动把通过程序制作出来的PPT手动复制到这个模板里面就可以了
""")
with open("Kmart Buy Trip Template.pptx", "rb") as f:
    btn = st.download_button(
        label="📥 点击下载 Kmart Buy Trip Template.pptx 文件",
        data=f,
        file_name="Kmart Buy Trip Template.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )